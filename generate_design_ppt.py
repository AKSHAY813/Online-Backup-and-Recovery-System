from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_design_presentation():
    prs = Presentation()

    # Define some colors
    BLUE = RGBColor(0, 82, 161)
    DARK = RGBColor(15, 23, 42)
    WHITE = RGBColor(255, 255, 255)
    SLATE = RGBColor(100, 116, 139)
    EMERALD = RGBColor(16, 185, 129)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Slide 1: Title - The Design Vision
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK)
    if os.path.exists("presentation_assets/landing_page.png"):
        slide.shapes.add_picture("presentation_assets/landing_page.png", 0, 0, width=prs.slide_width, height=prs.slide_height)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(4), prs.slide_width - Inches(1), Inches(2))
    tf = title.text_frame
    p = tf.add_paragraph()
    p.text = "NEXT-GEN UI/UX"
    p.font.bold = True
    p.font.size = Pt(72)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), prs.slide_width - Inches(1), Inches(1))
    p = subtitle.text_frame.add_paragraph()
    p.text = "Designing the Future of Data Recovery"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Design Philosophy
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, DARK)
    slide.shapes.title.text = "Design Philosophy: Glassmorphism"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    body = slide.placeholders[1]
    tf = body.text_frame
    for point in ["Visual Depth: Layers of frosted glass for hierarchy", "Minimalism: Removing friction from data management", "Feedback: Micro-animations for system health", "Accessibility: High-contrast typography (Outfit/Inter)"]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.color.rgb = WHITE

    # Slide 3: The Component Kit (UI Visuals)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK)
    if os.path.exists("presentation_assets/ui_kit.png"):
        slide.shapes.add_picture("presentation_assets/ui_kit.png", Inches(4), Inches(1), height=Inches(5.5))
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4), Inches(1))
    p = title.text_frame.add_paragraph()
    p.text = "Atomic Components"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = WHITE

    desc = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(3.5), Inches(3))
    tf = desc.text_frame
    for point in ["Neural Gauges", "Shard List Cards", "Dynamic Buttons", "Status Badges"]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.color.rgb = SLATE
        p.font.size = Pt(22)

    # Slide 4: Responsive Excellence (Mobile)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK)
    if os.path.exists("presentation_assets/mobile_ui.png"):
        slide.shapes.add_picture("presentation_assets/mobile_ui.png", Inches(0.5), Inches(1), height=Inches(5.5))
    
    title = slide.shapes.add_textbox(Inches(5.5), Inches(1), Inches(4), Inches(1))
    p = title.text_frame.add_paragraph()
    p.text = "Mobile Nodes"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = EMERALD

    desc = slide.shapes.add_textbox(Inches(5.5), Inches(2.5), Inches(4), Inches(3))
    tf = desc.text_frame
    for point in ["Capacitor-Powered", "Touch-Optimized Sharding", "Instant Cloud Sync", "Native-Feel Interactions"]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.color.rgb = WHITE
        p.font.size = Pt(22)

    # Slide 5: Color Palette & Typography
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, DARK)
    slide.shapes.title.text = "The Visual Identity"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    p = tf.add_paragraph()
    p.text = "Colors: Vault Blue, Slate Dark, Emerald Green"
    p.font.color.rgb = BLUE
    p = tf.add_paragraph()
    p.text = "Typography: 'Outfit' for headings, 'Inter' for body"
    p.font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "Motion: Neural-Pulse & Float animations"
    p.font.color.rgb = WHITE

    # Slide 6: User Flow (The UX)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, DARK)
    slide.shapes.title.text = "The User Journey"
    
    body = slide.placeholders[1]
    tf = body.text_frame
    for point in ["One-Click Sync: Automated mesh sharding", "Live Monitoring: Real-time efficiency gauges", "Temporal Recovery: Selecting snapshots from timeline", "Cognitive Parity: Secure reassembly of data"]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.color.rgb = WHITE

    # Slide 7: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BLUE)
    title = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    p = title.text_frame.add_paragraph()
    p.text = "WHERE SECURITY MEETS DESIGN"
    p.font.bold = True
    p.font.size = Pt(48)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    prs.save('Neural_Mesh_Design_Showcase.pptx')
    print("Design Showcase PPT created successfully!")

if __name__ == "__main__":
    create_design_presentation()
