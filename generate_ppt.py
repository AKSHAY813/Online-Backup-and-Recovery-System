from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()

    # Define some colors
    BLUE = RGBColor(0, 82, 161)
    DARK = RGBColor(15, 23, 42)
    WHITE = RGBColor(255, 255, 255)
    EMERALD = RGBColor(16, 185, 129)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Slide 1: Title
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DARK)
    
    if os.path.exists("presentation_assets/bg_title.png"):
        slide.shapes.add_picture("presentation_assets/bg_title.png", 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Add semi-transparent overlay (approximation with a shape)
    # title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    title = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), prs.slide_width - Inches(1), Inches(2))
    tf = title.text_frame
    p = tf.add_paragraph()
    p.text = "NEURAL MESH"
    p.font.bold = True
    p.font.size = Pt(80)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), prs.slide_width - Inches(1), Inches(1))
    stf = subtitle.text_frame
    sp = stf.add_paragraph()
    sp.text = "Cognitive Vault & Sharding Infrastructure"
    sp.font.size = Pt(24)
    sp.font.color.rgb = WHITE
    sp.alignment = PP_ALIGN.CENTER

    # Slide 2: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, DARK)
    title = slide.shapes.title
    title.text = "The Vulnerability Crisis"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Why Traditional Backups Fail"
    for point in ["Centralized Vulnerability (Single Point of Failure)", "Passive Storage (No Intelligence)", "Latency & Bottlenecks in Large Recovery"]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.level = 1
        p.font.color.rgb = WHITE

    # Slide 3: The Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, DARK)
    title = slide.shapes.title
    title.text = "Neural Mesh Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "A Paradigm Shift in Data Integrity"
    for point in ["P2P Sharding: Distributed Intelligence", "Cognitive Parity: AI-Driven Checks", "Mesh Guard: Real-time Protection Core"]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.level = 1
        p.font.color.rgb = WHITE

    # Slide 4: Key Feature - Dashboard
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK)
    if os.path.exists("presentation_assets/visual_dashboard.png"):
        slide.shapes.add_picture("presentation_assets/visual_dashboard.png", Inches(5), Inches(1), height=Inches(5.5))
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4), Inches(1))
    p = title.text_frame.add_paragraph()
    p.text = "Neural Monitor"
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = BLUE

    desc = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4), Inches(3))
    tf = desc.text_frame
    tf.word_wrap = True
    for point in ["Active Asset Cycle Analysis", "Predictive Storage Insights", "Neural Compression Gauges"]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.color.rgb = WHITE
        p.font.size = Pt(20)

    # Slide 5: Secure Recovery
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK)
    if os.path.exists("presentation_assets/visual_security.png"):
        slide.shapes.add_picture("presentation_assets/visual_security.png", Inches(0.5), Inches(1), height=Inches(5.5))
    
    title = slide.shapes.add_textbox(Inches(5.5), Inches(1), Inches(4), Inches(1))
    p = title.text_frame.add_paragraph()
    p.text = "Vault Recovery"
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = EMERALD

    desc = slide.shapes.add_textbox(Inches(5.5), Inches(2.5), Inches(4), Inches(3))
    tf = desc.text_frame
    tf.word_wrap = True
    for point in ["Temporal Shard Access", "Cognitive Parity Checksum", "Point-in-Time Recovery"]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.color.rgb = WHITE
        p.font.size = Pt(20)

    # Slide 6: Tech Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, DARK)
    title = slide.shapes.title
    title.text = "Technological Foundation"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    
    body = slide.placeholders[1]
    tf = body.text_frame
    for point in ["React & Vite: Performance Core", "Firebase: Real-time Mesh Layer", "Capacitor: Mobile Node Integration", "Tailwind CSS: Modern Aesthetic"]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.color.rgb = WHITE
        p.font.size = Pt(24)

    # Slide 7: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BLUE)
    title = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    p = title.text_frame.add_paragraph()
    p.text = "THE NEW STANDARD IN DATA INTEGRITY"
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    prs.save('Neural_Mesh_Presentation.pptx')
    print("Presentation created successfully as 'Neural_Mesh_Presentation.pptx'")

if __name__ == "__main__":
    create_presentation()
